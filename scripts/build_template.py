"""Rebuild graduation-template.pptx with an embedded placeholder image.

v2.5.1 fix: the template previously had a decorative rectangle where the
child's photo would go, but NO embedded image. Our v2.5.0 per-child
photo-swap logic needs an existing image relationship to rewrite. This
script adds a placeholder JPEG at the exact position of the photo-frame
rectangle so the swap has something to replace.

Run: python scripts/build_template.py
Output: src-tauri/resources/templates/graduation-template.pptx
"""
import io
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Emu

REPO = Path(__file__).resolve().parent.parent
TPL = REPO / "src-tauri" / "resources" / "templates" / "graduation-template.pptx"

# Rectangle 6 position + size from the current template's slide1.xml
PHOTO_X = Emu(6675120)
PHOTO_Y = Emu(822960)
PHOTO_CX = Emu(4846320)
PHOTO_CY = Emu(5212080)


def make_placeholder_jpg() -> bytes:
    """A warm cream headshot placeholder with a silhouette hint. The
    exact appearance doesn't matter — this image is REPLACED per child
    at render time. It just needs to exist so pptx.rs has an image
    relationship to swap."""
    w, h = 800, 860  # matches ~4846:5212 aspect
    img = Image.new("RGB", (w, h), (240, 229, 213))  # F0E5D5 (template cream)
    draw = ImageDraw.Draw(img)
    # Subtle silhouette (head + shoulders) in a warmer tone
    silhouette = (198, 176, 148)
    # Head (circle)
    hx, hy, hr = w // 2, int(h * 0.34), int(h * 0.14)
    draw.ellipse((hx - hr, hy - hr, hx + hr, hy + hr), fill=silhouette)
    # Shoulders (rounded)
    sy = int(h * 0.60)
    sx1, sx2 = int(w * 0.20), int(w * 0.80)
    draw.rounded_rectangle(
        (sx1, sy, sx2, h + 200),
        radius=int(w * 0.25),
        fill=silhouette,
    )
    # Save as JPEG
    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=85, optimize=True)
    return buf.getvalue()


def main() -> None:
    if not TPL.exists():
        raise SystemExit(f"missing template: {TPL}")
    pres = Presentation(str(TPL))
    slide = pres.slides[0]

    # Remove the "Drop {{Name}}'s photo here" text hint (TextBox 7 in
    # the raw XML). Iterate a copy so we can mutate the shape tree.
    to_remove = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            if "photo here" in text.lower() or "paste over this rectangle" in text.lower():
                to_remove.append(shape)
    for shape in to_remove:
        sp = shape._element
        sp.getparent().remove(sp)

    # Add the placeholder image at the photo-frame position.
    img_bytes = make_placeholder_jpg()
    pic = slide.shapes.add_picture(
        io.BytesIO(img_bytes),
        left=PHOTO_X,
        top=PHOTO_Y,
        width=PHOTO_CX,
        height=PHOTO_CY,
    )
    print(f"Added placeholder image: {pic.name} at "
          f"({PHOTO_X}, {PHOTO_Y}) {PHOTO_CX}x{PHOTO_CY} EMU")

    pres.save(str(TPL))
    size = TPL.stat().st_size
    print(f"Wrote {TPL} ({size:,} bytes)")


if __name__ == "__main__":
    main()
