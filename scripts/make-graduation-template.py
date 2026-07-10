"""Generate a designed graduation slide template for Echelon Receipts.

Produces a single-slide .pptx with placeholders {{Name}}, {{Note}}, {{Year}}.
The Graduation Day templater in src-tauri/src/graduation/pptx.rs clones this
slide once per graduating student and substitutes the placeholders.

Design register: warm, editorial, restrained — daycare / early-childhood
context. Not corporate SaaS, not cartoon-kids-clip-art. Palette: cream +
deep terracotta accent + warm ink. One serif + one geometric sans (contrast
axis, not two similar).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# 16:9 canvas at 13.333in x 7.5in — the modern PowerPoint default.
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Palette (OKLCH-informed hex):
CREAM = RGBColor(0xFA, 0xF6, 0xEE)     # body bg
INK = RGBColor(0x2A, 0x1F, 0x1A)       # near-black warm
TERRA = RGBColor(0xB4, 0x53, 0x28)     # accent — deep terracotta
MUTED = RGBColor(0x6B, 0x59, 0x4E)     # secondary text
GOLD = RGBColor(0xC9, 0x99, 0x40)      # divider accent

blank_layout = prs.slide_layouts[6]  # blank
slide = prs.slides.add_slide(blank_layout)

# Full-bleed cream background.
bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
)
bg.fill.solid()
bg.fill.fore_color.rgb = CREAM
bg.line.fill.background()

def add_text(left, top, width, height, text, *, size, color, bold=False,
             italic=False, align=PP_ALIGN.LEFT, font="Georgia"):
    """Add a text box with the given properties. Returns the shape."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return tb

# ── Left column: kid identity ─────────────────────────────────────────
LEFT_X = Inches(0.75)
LEFT_W = Inches(6.0)

# Small kicker year — the ONLY numbered element on the deck. Justified
# because it names the actual sequence (graduating class of {{Year}}).
add_text(LEFT_X, Inches(0.9), LEFT_W, Inches(0.4),
         "GRADUATING · CLASS OF {{Year}}",
         size=12, color=MUTED, font="Helvetica", bold=True)

# Big name — serif for warmth. Georgia is on every macOS and Windows box.
add_text(LEFT_X, Inches(1.35), LEFT_W, Inches(2.0),
         "{{Name}}",
         size=72, color=INK, font="Georgia", bold=True)

# Terracotta divider under the name.
div = slide.shapes.add_connector(1, LEFT_X, Inches(3.6), LEFT_X + Inches(1.2), Inches(3.6))
div.line.color.rgb = TERRA
div.line.width = Pt(3)

# Teacher note. Serif italic — reads like handwritten prose, not UI copy.
add_text(LEFT_X, Inches(3.9), LEFT_W, Inches(3.0),
         "{{Note}}",
         size=18, color=INK, italic=True, font="Georgia")

# ── Right column: photo placeholder ───────────────────────────────────
RIGHT_X = Inches(7.3)
RIGHT_W = Inches(5.3)
RIGHT_H = Inches(5.7)

# Thin terracotta frame that a photo will land inside. The user can
# paste a real image directly on top of this rectangle after the deck is
# generated — that's the intended workflow for v2.3.0. Photo-substitution
# from a per-child folder is a future feature.
frame = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, RIGHT_X, Inches(0.9), RIGHT_W, RIGHT_H
)
frame.fill.solid()
frame.fill.fore_color.rgb = RGBColor(0xF0, 0xE5, 0xD5)  # slightly warmer cream
frame.line.color.rgb = TERRA
frame.line.width = Pt(1)

# Instructions inside the frame — cleared once the user pastes a real
# photo over the top. Two lines, muted.
guide = slide.shapes.add_textbox(RIGHT_X, Inches(3.3), RIGHT_W, Inches(1.0))
gtf = guide.text_frame
gtf.word_wrap = True
p1 = gtf.paragraphs[0]
p1.alignment = PP_ALIGN.CENTER
r1 = p1.add_run()
r1.text = "Drop {{Name}}'s photo here"
r1.font.name = "Helvetica"
r1.font.size = Pt(16)
r1.font.color.rgb = MUTED
r1.font.italic = True
p2 = gtf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run()
r2.text = "(paste over this rectangle after generating)"
r2.font.name = "Helvetica"
r2.font.size = Pt(11)
r2.font.color.rgb = MUTED
r2.font.italic = True

# ── Footer band ──────────────────────────────────────────────────────
footer_line = slide.shapes.add_connector(
    1, Inches(0.75), Inches(6.9), Inches(12.58), Inches(6.9)
)
footer_line.line.color.rgb = GOLD
footer_line.line.width = Pt(0.75)

add_text(Inches(0.75), Inches(7.0), Inches(8), Inches(0.35),
         "ECHELON DAYCARE  ·  GRADUATION CEREMONY  ·  {{Year}}",
         size=10, color=MUTED, font="Helvetica")

out = Path(r"C:\Users\alosing\Desktop\echelon-receipts\src-tauri\resources\templates\graduation-template.pptx")
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(out)
print(f"wrote {out} ({out.stat().st_size} bytes)")
